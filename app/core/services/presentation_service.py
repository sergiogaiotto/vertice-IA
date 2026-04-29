"""Use case: Galeria de Apresentações VIP.

Gera apresentações executivas elegantes a partir do conjunto de cards
(insights, gráficos, diagramas) presentes em uma funcionalidade. Suporta
preview/edição antes de salvar e exportação para DOCX e PPTX com identidade
visual da plataforma Vértice (vermelho-marca + Inter).
"""

from __future__ import annotations

import io
import json
import re
import uuid
from dataclasses import dataclass
from datetime import datetime
from typing import Any

from app.adapters.db.sqlite import connect


# Identidade visual Vértice (RGB)
VERT_BRAND_DARK    = (220, 38, 38)     # #DC2626 — vermelho-marca
VERT_BRAND_LIGHT   = (254, 226, 226)   # #FEE2E2 — vermelho claríssimo
VERT_NEUTRAL_DARK  = (28, 25, 23)      # #1C1917
VERT_NEUTRAL_MED   = (87, 83, 78)      # #57534E
VERT_NEUTRAL_LIGHT = (250, 250, 249)   # #FAFAF9
VERT_AI_DARK       = (60, 52, 137)     # #3C3489


@dataclass
class Presentation:
    id: str
    title: str
    subtitle: str
    feature: str
    case_number: str
    sections: list[dict]   # [{title, body, source_card_uid?}]
    insights: list[dict]   # [{type, content}]
    visuals: list[dict]    # [{title, type, image_b64, caption, source_card_uid}]
    chat_history: list[dict]
    created_by_user: str
    created_by_id: str | None
    created_at: datetime
    updated_at: datetime
    cost_estimated: float
    tokens_input: int
    tokens_output: int
    model_used: str


def _row_to_presentation(row) -> Presentation:
    return Presentation(
        id=row[0],
        title=row[1],
        subtitle=row[2] or "",
        feature=row[3] or "",
        case_number=row[4] or "",
        sections=json.loads(row[5]) if row[5] else [],
        insights=json.loads(row[6]) if row[6] else [],
        chat_history=json.loads(row[7]) if row[7] else [],
        created_by_user=row[8] or "",
        created_by_id=row[9],
        created_at=datetime.fromisoformat(row[10]) if isinstance(row[10], str) else (row[10] or datetime.utcnow()),
        updated_at=datetime.fromisoformat(row[11]) if isinstance(row[11], str) else (row[11] or datetime.utcnow()),
        cost_estimated=float(row[12] or 0),
        tokens_input=int(row[13] or 0),
        tokens_output=int(row[14] or 0),
        model_used=row[15] or "",
        visuals=json.loads(row[16]) if (len(row) > 16 and row[16]) else [],
    )


_SELECT = (
    "SELECT id, title, subtitle, feature, case_number, sections, insights, "
    "chat_history, created_by_user, created_by_id, created_at, updated_at, "
    "cost_estimated, tokens_input, tokens_output, model_used, visuals "
    "FROM presentations"
)


class PresentationService:

    # ---------- persistência ----------

    async def save(self, p: Presentation) -> Presentation:
        async with connect() as db:
            await db.execute(
                "INSERT INTO presentations (id, title, subtitle, feature, case_number, "
                "sections, insights, chat_history, created_by_user, created_by_id, "
                "cost_estimated, tokens_input, tokens_output, model_used, visuals) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                (
                    p.id, p.title, p.subtitle, p.feature, p.case_number,
                    json.dumps(p.sections, ensure_ascii=False),
                    json.dumps(p.insights, ensure_ascii=False),
                    json.dumps(p.chat_history, ensure_ascii=False),
                    p.created_by_user, p.created_by_id,
                    p.cost_estimated, p.tokens_input, p.tokens_output, p.model_used,
                    json.dumps(p.visuals, ensure_ascii=False),
                ),
            )
            await db.commit()
        return p

    async def update(self, p: Presentation) -> Presentation:
        async with connect() as db:
            await db.execute(
                "UPDATE presentations SET title=?, subtitle=?, sections=?, insights=?, "
                "chat_history=?, visuals=?, updated_at=CURRENT_TIMESTAMP WHERE id=?",
                (
                    p.title, p.subtitle,
                    json.dumps(p.sections, ensure_ascii=False),
                    json.dumps(p.insights, ensure_ascii=False),
                    json.dumps(p.chat_history, ensure_ascii=False),
                    json.dumps(p.visuals, ensure_ascii=False),
                    p.id,
                ),
            )
            await db.commit()
        return p

    async def get(self, presentation_id: str) -> Presentation | None:
        async with connect() as db:
            cur = await db.execute(f"{_SELECT} WHERE id = ?", (presentation_id,))
            row = await cur.fetchone()
            return _row_to_presentation(row) if row else None

    async def list_all(self, limit: int = 100, q: str = "", feature: str = "") -> list[Presentation]:
        where = []
        params: list = []
        if q:
            where.append("(title LIKE ? OR subtitle LIKE ?)")
            params.extend([f"%{q}%", f"%{q}%"])
        if feature:
            where.append("feature = ?"); params.append(feature)
        clause = (" WHERE " + " AND ".join(where)) if where else ""
        async with connect() as db:
            cur = await db.execute(
                f"{_SELECT}{clause} ORDER BY created_at DESC LIMIT ?",
                (*params, limit),
            )
            return [_row_to_presentation(r) for r in await cur.fetchall()]

    async def delete(self, presentation_id: str) -> None:
        async with connect() as db:
            await db.execute("DELETE FROM presentations WHERE id = ?", (presentation_id,))
            await db.commit()

    async def stats(self) -> dict:
        async with connect() as db:
            cur = await db.execute(
                "SELECT COUNT(*), SUM(cost_estimated), SUM(tokens_input + tokens_output) FROM presentations"
            )
            row = await cur.fetchone()
            return {
                "total": int(row[0] or 0),
                "total_cost": float(row[1] or 0),
                "total_tokens": int(row[2] or 0),
            }

    # ---------- factory ----------

    @staticmethod
    def new(title: str, subtitle: str, feature: str, case_number: str,
            sections: list[dict], insights: list[dict],
            user: str, user_id: str | None,
            visuals: list[dict] | None = None,
            tokens_in: int = 0, tokens_out: int = 0, cost: float = 0.0,
            model: str = "") -> Presentation:
        return Presentation(
            id=uuid.uuid4().hex,
            title=title, subtitle=subtitle,
            feature=feature, case_number=case_number,
            sections=sections, insights=insights, chat_history=[],
            visuals=visuals or [],
            created_by_user=user, created_by_id=user_id,
            created_at=datetime.utcnow(), updated_at=datetime.utcnow(),
            cost_estimated=cost, tokens_input=tokens_in, tokens_output=tokens_out,
            model_used=model,
        )

    # ---------- exports ----------

    def _decode_visual(self, vis: dict) -> bytes | None:
        """Extrai bytes da imagem b64 (data:image/png;base64,XYZ ou XYZ puro)."""
        import base64
        b64 = (vis.get("image_b64") or "").strip()
        if not b64:
            return None
        if b64.startswith("data:"):
            comma = b64.find(",")
            if comma >= 0:
                b64 = b64[comma + 1:]
        try:
            return base64.b64decode(b64)
        except Exception:
            return None

    def export_docx(self, p: Presentation) -> bytes:
        """Export DOCX elegante com identidade Vértice."""
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_ALIGN_VERTICAL
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()
        # margens generosas para visual VIP
        for section in doc.sections:
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(2.8)
            section.right_margin = Cm(2.8)

        # ----- Capa -----
        # marca verde no topo
        marca = doc.add_paragraph()
        marca_run = marca.add_run("VÉRTICE")
        marca_run.font.name = "Inter"
        marca_run.font.size = Pt(11)
        marca_run.font.bold = True
        marca_run.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
        marca.paragraph_format.space_after = Pt(0)

        marca_sub = doc.add_paragraph()
        marca_sub_run = marca_sub.add_run("APRESENTAÇÃO EXECUTIVA")
        marca_sub_run.font.name = "Inter"
        marca_sub_run.font.size = Pt(8)
        marca_sub_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
        marca_sub.paragraph_format.space_after = Pt(60)

        # Título
        title_p = doc.add_paragraph()
        title_run = title_p.add_run(p.title)
        title_run.font.name = "Inter"
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
        title_p.paragraph_format.space_after = Pt(8)

        if p.subtitle:
            sub_p = doc.add_paragraph()
            sub_run = sub_p.add_run(p.subtitle)
            sub_run.font.name = "Inter"
            sub_run.font.size = Pt(14)
            sub_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
            sub_p.paragraph_format.space_after = Pt(40)

        # Metadados em rodapé da capa
        meta = doc.add_paragraph()
        meta_run = meta.add_run(
            f"{p.created_at.strftime('%d/%m/%Y')}"
            + (f"  ·  {p.feature}" if p.feature else "")
            + (f"  ·  caso {p.case_number}" if p.case_number else "")
        )
        meta_run.font.name = "Inter"
        meta_run.font.size = Pt(9)
        meta_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
        meta.paragraph_format.space_after = Pt(60)

        # ----- Insights (Executive summary) -----
        if p.insights:
            doc.add_page_break()
            h = doc.add_paragraph()
            h_run = h.add_run("INSIGHTS PRINCIPAIS")
            h_run.font.name = "Inter"
            h_run.font.size = Pt(10)
            h_run.font.bold = True
            h_run.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
            h.paragraph_format.space_after = Pt(16)

            for ins in p.insights:
                self._add_insight_block(doc, ins)

        # ----- Seções -----
        for sec in p.sections:
            doc.add_page_break()
            sec_h = doc.add_paragraph()
            sec_h_run = sec_h.add_run(sec.get("title", "").upper())
            sec_h_run.font.name = "Inter"
            sec_h_run.font.size = Pt(10)
            sec_h_run.font.bold = True
            sec_h_run.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
            sec_h.paragraph_format.space_after = Pt(4)

            sec_t = doc.add_paragraph()
            sec_t_run = sec_t.add_run(sec.get("title", ""))
            sec_t_run.font.name = "Inter"
            sec_t_run.font.size = Pt(20)
            sec_t_run.font.bold = True
            sec_t_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
            sec_t.paragraph_format.space_after = Pt(20)

            self._render_markdown_to_docx(doc, sec.get("body", ""))

        # ----- Visuais (gráficos e diagramas) — uma página por visual -----
        for vis in p.visuals:
            img_bytes = self._decode_visual(vis)
            if not img_bytes:
                continue
            doc.add_page_break()

            # label de tipo
            type_p = doc.add_paragraph()
            type_run = type_p.add_run((vis.get("type", "VISUAL")).upper())
            type_run.font.name = "Inter"
            type_run.font.size = Pt(10)
            type_run.font.bold = True
            type_run.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
            type_p.paragraph_format.space_after = Pt(4)

            # título
            title_p = doc.add_paragraph()
            title_run = title_p.add_run(vis.get("title", "Visual"))
            title_run.font.name = "Inter"
            title_run.font.size = Pt(20)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
            title_p.paragraph_format.space_after = Pt(16)

            # imagem (largura útil ~15.5cm com margens)
            try:
                img_p = doc.add_paragraph()
                img_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                img_run = img_p.add_run()
                img_run.add_picture(io.BytesIO(img_bytes), width=Cm(15.5))
            except Exception:
                err_p = doc.add_paragraph()
                err_run = err_p.add_run("[imagem não pôde ser inserida]")
                err_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
                err_run.italic = True

            # caption
            if vis.get("caption"):
                cap_p = doc.add_paragraph()
                cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cap_run = cap_p.add_run(vis["caption"])
                cap_run.font.name = "Inter"
                cap_run.font.size = Pt(10)
                cap_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
                cap_run.italic = True
                cap_p.paragraph_format.space_before = Pt(8)

        # rodapé final
        doc.add_paragraph()
        footer = doc.add_paragraph()
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_run = footer.add_run(
            f"Gerado em Vértice  ·  {p.created_at.strftime('%d/%m/%Y %H:%M')}"
            + (f"  ·  por {p.created_by_user}" if p.created_by_user else "")
        )
        footer_run.font.name = "Inter"
        footer_run.font.size = Pt(8)
        footer_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _add_insight_block(self, doc, insight: dict) -> None:
        """Bloco de insight: barra lateral verde + texto."""
        from docx.shared import Pt, RGBColor, Cm
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        ins_type = insight.get("type", "").upper()
        content = insight.get("content", "")

        # tipo (ex: "OPORTUNIDADE", "RISCO", "DESTAQUE")
        if ins_type:
            type_p = doc.add_paragraph()
            type_run = type_p.add_run(ins_type)
            type_run.font.name = "Inter"
            type_run.font.size = Pt(8)
            type_run.font.bold = True
            type_run.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
            type_p.paragraph_format.space_after = Pt(2)

        ins_p = doc.add_paragraph()
        ins_run = ins_p.add_run(content)
        ins_run.font.name = "Inter"
        ins_run.font.size = Pt(13)
        ins_run.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
        ins_p.paragraph_format.space_after = Pt(14)
        ins_p.paragraph_format.left_indent = Cm(0.5)

    def _render_markdown_to_docx(self, doc, md: str) -> None:
        """Renderiza markdown simples no doc (h2/h3/listas/parágrafos/código)."""
        from docx.shared import Pt, RGBColor, Cm

        for raw_line in (md or "").split("\n"):
            line = raw_line.rstrip()
            if not line.strip():
                continue

            if line.startswith("### "):
                p = doc.add_paragraph()
                r = p.add_run(line[4:].strip())
                r.font.name = "Inter"
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
                p.paragraph_format.space_before = Pt(12)
                p.paragraph_format.space_after = Pt(6)
            elif line.startswith("## "):
                p = doc.add_paragraph()
                r = p.add_run(line[3:].strip())
                r.font.name = "Inter"
                r.font.size = Pt(16)
                r.font.bold = True
                r.font.color.rgb = RGBColor(*VERT_BRAND_DARK)
                p.paragraph_format.space_before = Pt(16)
                p.paragraph_format.space_after = Pt(8)
            elif line.startswith("- ") or line.startswith("* "):
                p = doc.add_paragraph(style="List Bullet")
                r = p.add_run(line[2:].strip())
                r.font.name = "Inter"
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
            elif re.match(r"^\d+\.\s", line):
                p = doc.add_paragraph(style="List Number")
                r = p.add_run(re.sub(r"^\d+\.\s", "", line))
                r.font.name = "Inter"
                r.font.size = Pt(11)
                r.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)
            elif line.startswith("> "):
                p = doc.add_paragraph()
                r = p.add_run(line[2:].strip())
                r.font.name = "Inter"
                r.font.size = Pt(11)
                r.italic = True
                r.font.color.rgb = RGBColor(*VERT_NEUTRAL_MED)
                p.paragraph_format.left_indent = Cm(1)
            else:
                # parágrafo normal — processa **bold** inline básico
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(6)
                self._add_runs_with_bold(p, line)

    def _add_runs_with_bold(self, paragraph, text: str) -> None:
        """Suporta **negrito** dentro de parágrafo normal."""
        from docx.shared import Pt, RGBColor

        parts = re.split(r"(\*\*[^*]+\*\*)", text)
        for part in parts:
            if not part:
                continue
            is_bold = part.startswith("**") and part.endswith("**")
            content = part[2:-2] if is_bold else part
            r = paragraph.add_run(content)
            r.font.name = "Inter"
            r.font.size = Pt(11)
            r.font.bold = is_bold
            r.font.color.rgb = RGBColor(*VERT_NEUTRAL_DARK)

    def export_pptx(self, p: Presentation) -> bytes:
        """Export PPTX widescreen 16:9 com identidade Vértice."""
        from pptx import Presentation as PPTX
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = PPTX()
        # Widescreen 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank = prs.slide_layouts[6]

        # ===== Slide 1: Capa =====
        slide = prs.slides.add_slide(blank)
        # Fundo claro
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*VERT_NEUTRAL_LIGHT)
        bg.line.fill.background()

        # Faixa lateral verde (decoração elegante)
        side = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height
        )
        side.fill.solid()
        side.fill.fore_color.rgb = RGBColor(*VERT_BRAND_DARK)
        side.line.fill.background()

        # Marca topo
        marca = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(8), Inches(0.4))
        self._set_text(marca, "VÉRTICE", Pt(11), bold=True, color=VERT_BRAND_DARK, font="Inter")
        marca_sub = slide.shapes.add_textbox(Inches(0.9), Inches(1.0), Inches(8), Inches(0.4))
        self._set_text(marca_sub, "APRESENTAÇÃO EXECUTIVA", Pt(8), color=VERT_NEUTRAL_MED, font="Inter")

        # Título central
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(2))
        self._set_text(title_box, p.title, Pt(44), bold=True, color=VERT_NEUTRAL_DARK, font="Inter")

        if p.subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(11.5), Inches(1))
            self._set_text(sub_box, p.subtitle, Pt(20), color=VERT_NEUTRAL_MED, font="Inter")

        # Metadados rodapé
        meta_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.5))
        meta_text = f"{p.created_at.strftime('%d/%m/%Y')}"
        if p.feature: meta_text += f"  ·  {p.feature}"
        if p.case_number: meta_text += f"  ·  caso {p.case_number}"
        self._set_text(meta_box, meta_text, Pt(10), color=VERT_NEUTRAL_MED, font="Inter")

        # ===== Slide 2: Insights (executive summary) =====
        if p.insights:
            slide = prs.slides.add_slide(blank)
            self._add_slide_header(slide, prs, "INSIGHTS PRINCIPAIS")
            # cards de insights em 2 colunas
            n = len(p.insights)
            col_w = Inches(5.7)
            col_h = Inches(1.4)
            margin_top = Inches(1.6)
            margin_x = Inches(0.6)
            gap_x = Inches(0.3)
            gap_y = Inches(0.25)
            for i, ins in enumerate(p.insights[:8]):  # cap em 8 insights
                col = i % 2
                row = i // 2
                x = margin_x + col * (col_w + gap_x)
                y = margin_top + row * (col_h + gap_y)
                # caixa do card
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, col_h)
                card.adjustments[0] = 0.05
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card.line.color.rgb = RGBColor(*VERT_BRAND_LIGHT)
                card.line.width = Pt(0.75)
                # texto da label
                label_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), Inches(0.3))
                self._set_text(label_box, ins.get("type", "INSIGHT").upper(), Pt(9), bold=True, color=VERT_BRAND_DARK, font="Inter")
                # conteúdo
                content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.45), col_w - Inches(0.4), col_h - Inches(0.55))
                self._set_text(content_box, ins.get("content", ""), Pt(11), color=VERT_NEUTRAL_DARK, font="Inter")

        # ===== Slides de seção (1 por seção) =====
        for sec in p.sections:
            slide = prs.slides.add_slide(blank)
            self._add_slide_header(slide, prs, sec.get("title", "").upper() or "SEÇÃO")

            # Título
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.7))
            self._set_text(title_box, sec.get("title", ""), Pt(28), bold=True, color=VERT_NEUTRAL_DARK, font="Inter")

            # Corpo
            body_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(5.0))
            self._render_markdown_to_pptx(body_box, sec.get("body", ""))

        # ===== Slides de visuais (gráficos e diagramas) =====
        for vis in p.visuals:
            img_bytes = self._decode_visual(vis)
            if not img_bytes:
                continue
            slide = prs.slides.add_slide(blank)
            self._add_slide_header(slide, prs, (vis.get("type", "VISUAL")).upper())

            # Título
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.7))
            self._set_text(title_box, vis.get("title", "Visual"), Pt(26), bold=True, color=VERT_NEUTRAL_DARK, font="Inter")

            # Imagem centralizada — calcula tamanho proporcional para caber em ~Inches(11×4.5)
            try:
                img_buf = io.BytesIO(img_bytes)
                # adiciona imagem temporária para descobrir dimensões nativas
                tmp_pic = slide.shapes.add_picture(img_buf, Inches(0.6), Inches(2.0))
                native_w = tmp_pic.width
                native_h = tmp_pic.height
                # calcula scale para caber
                max_w = Inches(11.0); max_h = Inches(4.5)
                scale = min(max_w / native_w, max_h / native_h, 1.0)
                final_w = int(native_w * scale)
                final_h = int(native_h * scale)
                # remove temp e re-adiciona centralizado
                spt = tmp_pic._element
                spt.getparent().remove(spt)
                left = int((prs.slide_width - final_w) / 2)
                top = Inches(2.1) + int((Inches(4.5) - final_h) / 2)
                img_buf.seek(0)
                slide.shapes.add_picture(img_buf, left, top, width=final_w, height=final_h)
            except Exception:
                err_box = slide.shapes.add_textbox(Inches(0.6), Inches(3), Inches(12.1), Inches(0.5))
                self._set_text(err_box, "[imagem não pôde ser inserida]", Pt(11), color=VERT_NEUTRAL_MED, font="Inter")

            # Caption no rodapé
            if vis.get("caption"):
                cap_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
                self._set_text(cap_box, vis["caption"], Pt(11), color=VERT_NEUTRAL_MED, font="Inter")

        # ===== Slide final =====
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(*VERT_BRAND_DARK); bg.line.fill.background()
        thanks_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.5))
        self._set_text(thanks_box, "Obrigado.", Pt(54), bold=True, color=(255, 255, 255), font="Inter", align=PP_ALIGN.CENTER)
        info_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.5))
        info_text = f"Gerado em Vértice · {p.created_at.strftime('%d/%m/%Y %H:%M')}"
        if p.created_by_user: info_text += f" · por {p.created_by_user}"
        self._set_text(info_box, info_text, Pt(11), color=VERT_BRAND_LIGHT, font="Inter", align=PP_ALIGN.CENTER)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _add_slide_header(self, slide, prs, label: str) -> None:
        """Header padrão dos slides internos: marca + label."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        # linha decorativa verde no topo
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(*VERT_BRAND_DARK); line.line.fill.background()

        # marca topo direito
        marca = slide.shapes.add_textbox(prs.slide_width - Inches(2), Inches(0.25), Inches(1.7), Inches(0.4))
        self._set_text(marca, "VÉRTICE", Pt(8), bold=True, color=VERT_BRAND_DARK, font="Inter", align=2)  # right

        # label esquerdo
        label_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11), Inches(0.4))
        self._set_text(label_box, label, Pt(9), bold=True, color=VERT_BRAND_DARK, font="Inter")

    def _set_text(self, shape, text: str, size, bold: bool = False,
                  color: tuple = (0, 0, 0), font: str = "Inter", align=None) -> None:
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)

    def _render_markdown_to_pptx(self, shape, md: str) -> None:
        """Renderiza markdown simples num text_frame do PPTX."""
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True

        for raw in (md or "").split("\n"):
            line = raw.rstrip()
            if not line.strip():
                continue

            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            indent = 0
            size = Pt(13)
            bold = False
            color = VERT_NEUTRAL_DARK
            italic = False

            if line.startswith("### "):
                content = line[4:]
                size = Pt(15); bold = True
            elif line.startswith("## "):
                content = line[3:]
                size = Pt(18); bold = True; color = VERT_BRAND_DARK
            elif line.startswith("- ") or line.startswith("* "):
                content = "•  " + line[2:]
                indent = 1
            elif re.match(r"^\d+\.\s", line):
                content = line
                indent = 1
            elif line.startswith("> "):
                content = line[2:]
                italic = True; color = VERT_NEUTRAL_MED
            else:
                # bold inline
                content = re.sub(r"\*\*([^*]+)\*\*", r"\1", line)
                bold = "**" in line   # simplificação: se tem bold, marca todo o run

            p.level = indent
            run = p.add_run()
            run.text = content
            run.font.name = "Inter"
            run.font.size = size
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = RGBColor(*color)


_global = PresentationService()


def get_presentation_service() -> PresentationService:
    return _global
